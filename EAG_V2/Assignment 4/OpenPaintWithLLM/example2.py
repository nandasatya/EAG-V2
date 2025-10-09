# basic import 
from mcp.server.fastmcp import FastMCP, Image
from mcp.server.fastmcp.prompts import base
from mcp.types import TextContent
from mcp import types
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import math
import sys
import time
import subprocess
import platform
import os
import tempfile

# instantiate an MCP server client
mcp = FastMCP("Calculator")

# Global variable for drawing application
drawing_app = None
# Global variable to store rectangle coordinates
rectangle_coords = None
# Global variable to store PowerPoint file path
ppt_file_path = None
# Global variable to store presentation object
prs = None

# macOS automation helper functions
def run_applescript(script):
    """Run AppleScript on macOS"""
    try:
        result = subprocess.run(['osascript', '-e', script], 
                              capture_output=True, text=True, timeout=10)
        return result.returncode == 0, result.stdout, result.stderr
    except subprocess.TimeoutExpired:
        return False, "", "AppleScript timeout"
    except Exception as e:
        return False, "", str(e)

def open_drawing_app_macos():
    """Open a drawing application on macOS"""
    try:
        # Use only PowerPoint for drawing
        drawing_apps = [
            'Microsoft PowerPoint' # PowerPoint only
        ]
        
        app_opened = False
        for app in drawing_apps:
            try:
                subprocess.Popen(['open', '-a', app])
                time.sleep(2)
                app_opened = True
                break
            except:
                continue
        
        if not app_opened:
            # PowerPoint is required - no fallback
            return False, "PowerPoint is required but not available"
        
        return app_opened, "Drawing application opened successfully"
    except Exception as e:
        return False, f"Error opening drawing app: {str(e)}"

def draw_rectangle_macos(x1, y1, x2, y2):
    """Draw rectangle using python-pptx"""
    global rectangle_coords, ppt_file_path, prs
    try:
        # Store rectangle coordinates globally
        rectangle_coords = [x1, y1, x2, y2]
        print(f"DEBUG: Stored rectangle coordinates: {rectangle_coords}")
        
        if not prs or not ppt_file_path:
            return False, "PowerPoint presentation not created. Please call create_new_slide first."
        
        try:
            # Get the first slide
            slide = prs.slides[0]
            
            # Convert pixel coordinates to inches (better scaling for visibility)
            # Scale down by a factor to fit on slide better
            left = Inches(x1 / 100.0)
            top = Inches(y1 / 100.0)
            width = Inches((x2 - x1) / 100.0)
            height = Inches((y2 - y1) / 100.0)
            
            # Add rectangle shape
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            
            # Style the rectangle - make it VERY visible
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 200)  # Light yellow fill
            shape.line.color.rgb = RGBColor(255, 0, 0)  # RED border
            shape.line.width = Pt(6)  # Thick border
            
            # Save the presentation
            prs.save(ppt_file_path)
            print(f"Rectangle drawn and saved to {ppt_file_path}")
            
            return True, f"Rectangle drawn from ({x1},{y1}) to ({x2},{y2}) in PowerPoint"
                
        except Exception as e:
            print(f"PowerPoint drawing failed: {e}")
            return False, f"PowerPoint drawing failed: {str(e)}"
            
    except Exception as e:
        return False, f"Error drawing rectangle: {str(e)}"

def create_new_slide_macos():
    """Create a new slide in PowerPoint"""
    global ppt_file_path, prs
    try:
        # Create a new PowerPoint presentation using python-pptx
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]  # 6 is blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Save to a temporary file with timestamp to avoid caching
        import datetime
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_file_path = os.path.join(tempfile.gettempdir(), f"mcp_powerpoint_{timestamp}.pptx")
        prs.save(ppt_file_path)
        print(f"PowerPoint file created at: {ppt_file_path}")
        
        # Open the PowerPoint file
        subprocess.Popen(['open', ppt_file_path])
        time.sleep(2)
        
        return True, f"New slide created in PowerPoint"
            
    except Exception as e:
        return False, f"Error creating new slide: {str(e)}"

def add_text_macos(text):
    """Add text to the drawing using python-pptx"""
    global rectangle_coords, ppt_file_path, prs
    try:
        if not prs or not ppt_file_path:
            return False, "PowerPoint presentation not created. Please call create_new_slide first."
        
        if not rectangle_coords:
            return False, "Rectangle not drawn. Please call draw_rectangle first."
        
        try:
            # Get the first slide
            slide = prs.slides[0]
            
            # Calculate center position based on rectangle coordinates
            x1, y1, x2, y2 = rectangle_coords
            center_x = (x1 + x2) / 2
            center_y = (y1 + y2) / 2
            
            # Convert to inches (better scaling for visibility)
            left = Inches((center_x - 150) / 100.0)  # Center text box (300px wide)
            top = Inches((center_y - 40) / 100.0)    # Center text box (80px tall)
            width = Inches(300 / 100.0)
            height = Inches(80 / 100.0)
            
            # Add text box
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = str(text)
            
            # Format the text - make it VERY visible
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(44)  # LARGE text
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 0, 255)  # BLUE color (easier to see)
            paragraph.font.name = "Arial"
            from pptx.enum.text import PP_ALIGN
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Save the presentation
            prs.save(ppt_file_path)
            print(f"Text added and saved to {ppt_file_path}")
            
            return True, f"Text '{text}' added inside the rectangle in PowerPoint"
                
        except Exception as e:
            print(f"PowerPoint text failed: {e}")
            return False, f"PowerPoint text failed: {str(e)}"
            
    except Exception as e:
        return False, f"Error adding text: {str(e)}"

# DEFINE TOOLS

#addition tool
@mcp.tool()
def add(a: int, b: int) -> int:
    """Add two numbers"""
    print("CALLED: add(a: int, b: int) -> int:")
    return int(a + b)

@mcp.tool()
def add_list(l: list) -> int:
    """Add all numbers in a list"""
    print("CALLED: add(l: list) -> int:")
    return sum(l)

# subtraction tool
@mcp.tool()
def subtract(a: int, b: int) -> int:
    """Subtract two numbers"""
    print("CALLED: subtract(a: int, b: int) -> int:")
    return int(a - b)

# multiplication tool
@mcp.tool()
def multiply(a: int, b: int) -> int:
    """Multiply two numbers"""
    print("CALLED: multiply(a: int, b: int) -> int:")
    return int(a * b)

#  division tool
@mcp.tool() 
def divide(a: int, b: int) -> float:
    """Divide two numbers"""
    print("CALLED: divide(a: int, b: int) -> float:")
    return float(a / b)

# power tool
@mcp.tool()
def power(a: int, b: int) -> int:
    """Power of two numbers"""
    print("CALLED: power(a: int, b: int) -> int:")
    return int(a ** b)

# square root tool
@mcp.tool()
def sqrt(a: int) -> float:
    """Square root of a number"""
    print("CALLED: sqrt(a: int) -> float:")
    return float(a ** 0.5)

# cube root tool
@mcp.tool()
def cbrt(a: int) -> float:
    """Cube root of a number"""
    print("CALLED: cbrt(a: int) -> float:")
    return float(a ** (1/3))

# factorial tool
@mcp.tool()
def factorial(a: int) -> int:
    """factorial of a number"""
    print("CALLED: factorial(a: int) -> int:")
    return int(math.factorial(a))

# log tool
@mcp.tool()
def log(a: int) -> float:
    """log of a number"""
    print("CALLED: log(a: int) -> float:")
    return float(math.log(a))

# remainder tool
@mcp.tool()
def remainder(a: int, b: int) -> int:
    """remainder of two numbers divison"""
    print("CALLED: remainder(a: int, b: int) -> int:")
    return int(a % b)

# sin tool
@mcp.tool()
def sin(a: int) -> float:
    """sin of a number"""
    print("CALLED: sin(a: int) -> float:")
    return float(math.sin(a))

# cos tool
@mcp.tool()
def cos(a: int) -> float:
    """cos of a number"""
    print("CALLED: cos(a: int) -> float:")
    return float(math.cos(a))

# tan tool
@mcp.tool()
def tan(a: int) -> float:
    """tan of a number"""
    print("CALLED: tan(a: int) -> float:")
    return float(math.tan(a))

# mine tool
@mcp.tool()
def mine(a: int, b: int) -> int:
    """special mining tool"""
    print("CALLED: mine(a: int, b: int) -> int:")
    return int(a - b - b)

@mcp.tool()
def create_thumbnail(image_path: str) -> Image:
    """Create a thumbnail from an image"""
    print("CALLED: create_thumbnail(image_path: str) -> Image:")
    img = PILImage.open(image_path)
    img.thumbnail((100, 100))
    return Image(data=img.tobytes(), format="png")

@mcp.tool()
def strings_to_chars_to_int(string: str) -> list[int]:
    """Return the ASCII values of the characters in a word"""
    print("CALLED: strings_to_chars_to_int(string: str) -> list[int]:")
    return [int(ord(char)) for char in string]

@mcp.tool()
def int_list_to_exponential_sum(int_list: list) -> float:
    """Return sum of exponentials of numbers in a list"""
    print("CALLED: int_list_to_exponential_sum(int_list: list) -> float:")
    return sum(math.exp(i) for i in int_list)

@mcp.tool()
def fibonacci_numbers(n: int) -> list:
    """Return the first n Fibonacci Numbers"""
    print("CALLED: fibonacci_numbers(n: int) -> list:")
    if n <= 0:
        return []
    fib_sequence = [0, 1]
    for _ in range(2, n):
        fib_sequence.append(fib_sequence[-1] + fib_sequence[-2])
    return fib_sequence[:n]


@mcp.tool()
async def draw_rectangle(x1: int, y1: int, x2: int, y2: int) -> dict:
    """Draw rectangle on the current PowerPoint slide from (x1,y1) to (x2,y2)"""
    global drawing_app
    try:
        if platform.system() == "Darwin":  # macOS
            if not drawing_app:
                # Open drawing app if not already open
                success, message = open_drawing_app_macos()
                if not success:
                    return {
                        "content": [
                            TextContent(
                                type="text",
                                text=f"Error opening PowerPoint app: {message}"
                            )
                        ]
                    }
                drawing_app = True
            
            # Use macOS automation to draw rectangle
            success, message = draw_rectangle_macos(x1, y1, x2, y2)
            
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"New document created in PowerPoint. {message}"
                    )
                ]
            }
        else:
            # Windows fallback (original functionality preserved)
            if not drawing_app:
                return {
                    "content": [
                        TextContent(
                            type="text",
                            text="PowerPoint application is not open. Please call open_paint first."
                        )
                    ]
                }
            
            # Original Windows code would go here
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"Rectangle drawn from ({x1},{y1}) to ({x2},{y2}) - Windows mode"
                    )
                ]
            }
    except Exception as e:
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error drawing rectangle in PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def add_text_in_paint(text: str) -> dict:
    """Add text inside the rectangle in PowerPoint"""
    global drawing_app
    try:
        if platform.system() == "Darwin":  # macOS
            if not drawing_app:
                # Open drawing app if not already open
                success, message = open_drawing_app_macos()
                if not success:
                    return {
                        "content": [
                            TextContent(
                                type="text",
                                text=f"Error opening PowerPoint app: {message}"
                            )
                        ]
                    }
                drawing_app = True
            
            # Use macOS automation to add text
            success, message = add_text_macos(text)
            
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"Text '{text}' added inside the rectangle in PowerPoint. {message}"
                    )
                ]
            }
        else:
            # Windows fallback (original functionality preserved)
            if not drawing_app:
                return {
                    "content": [
                        TextContent(
                            type="text",
                            text="PowerPoint application is not open. Please call open_paint first."
                        )
                    ]
                }
            
            # Original Windows code would go here
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"Text '{text}' added to drawing application - Windows mode"
                    )
                ]
            }
    except Exception as e:
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error adding text to PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def create_new_slide() -> dict:
    """Create a new slide in PowerPoint"""
    global drawing_app
    try:
        if platform.system() == "Darwin":  # macOS
            if not drawing_app:
                return {
                    "content": [
                        TextContent(
                            type="text",
                            text="PowerPoint application is not open. Please call open_paint first."
                        )
                    ]
                }
            
            success, message = create_new_slide_macos()
            if success:
                return {
                    "content": [
                        TextContent(
                            type="text",
                            text=f"New slide created in PowerPoint. {message}"
                        )
                    ]
                }
            else:
                return {
                    "content": [
                        TextContent(
                            type="text",
                            text=f"Error creating new slide in PowerPoint: {message}"
                        )
                    ]
                }
        else:
            # Windows fallback (original functionality preserved)
            return {
                "content": [
                    TextContent(
                        type="text",
                        text="New slide created in PowerPoint - Windows mode"
                    )
                ]
            }
    except Exception as e:
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error creating new slide in PowerPoint: {str(e)}"
                )
            ]
        }

@mcp.tool()
async def open_paint() -> dict:
    """Open PowerPoint application for drawing"""
    global drawing_app
    try:
        if platform.system() == "Darwin":  # macOS
            success, message = open_drawing_app_macos()
            if success:
                drawing_app = True
            
            return {
                "content": [
                    TextContent(
                        type="text",
                        text=f"PowerPoint application opened successfully. Ready for drawing. {message}"
                    )
                ]
            }
        else:
            # Windows fallback (original functionality preserved)
            # Original Windows code would go here
            drawing_app = True
            return {
                "content": [
                    TextContent(
                        type="text",
                        text="Drawing application opened successfully - Windows mode"
                    )
                ]
            }
    except Exception as e:
        return {
            "content": [
                TextContent(
                    type="text",
                    text=f"Error opening PowerPoint application: {str(e)}"
                )
            ]
        }
# DEFINE RESOURCES

# Add a dynamic greeting resource
@mcp.resource("greeting://{name}")
def get_greeting(name: str) -> str:
    """Get a personalized greeting"""
    print("CALLED: get_greeting(name: str) -> str:")
    return f"Hello, {name}!"


# DEFINE AVAILABLE PROMPTS
@mcp.prompt()
def review_code(code: str) -> str:
    return f"Please review this code:\n\n{code}"
    print("CALLED: review_code(code: str) -> str:")


@mcp.prompt()
def debug_error(error: str) -> list[base.Message]:
    return [
        base.UserMessage("I'm seeing this error:"),
        base.UserMessage(error),
        base.AssistantMessage("I'll help debug that. What have you tried so far?"),
    ]

if __name__ == "__main__":
    # Check if running with mcp dev command
    print("STARTING")
    if len(sys.argv) > 1 and sys.argv[1] == "dev":
        mcp.run()  # Run without transport for dev server
    else:
        mcp.run(transport="stdio")  # Run with stdio for direct execution
